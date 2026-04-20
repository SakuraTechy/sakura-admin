from __future__ import annotations

import base64
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "ai-platform-leadership-deck.md"
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "ai-platform-leadership-deck.pptx"

PURPLE = RGBColor(122, 63, 166)
PURPLE_LIGHT = RGBColor(242, 233, 250)
TEXT = RGBColor(51, 51, 51)
GRAY = RGBColor(108, 108, 108)
BLUE = RGBColor(36, 132, 219)
BG = RGBColor(252, 252, 252)


@dataclass
class SlideContent:
    title: str
    points: List[str]
    script: str
    mermaid: str | None


def parse_slides(text: str) -> List[SlideContent]:
    pattern = re.compile(r"^## (第 .*?)$", re.MULTILINE)
    matches = list(pattern.finditer(text))
    slides: List[SlideContent] = []

    for idx, match in enumerate(matches):
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        block = text[start:end].strip()
        title = match.group(1).strip()

        mermaid_match = re.search(r"```mermaid\s*(.*?)```", block, re.S)
        mermaid = mermaid_match.group(1).strip() if mermaid_match else None
        content_wo_mermaid = re.sub(r"```mermaid\s*.*?```", "", block, flags=re.S).strip()

        script_match = re.search(r"### 讲解话术\s*(.*)$", content_wo_mermaid, re.S)
        script = ""
        if script_match:
            script = re.sub(r"\n+", " ", script_match.group(1)).strip()
            content_wo_mermaid = content_wo_mermaid[: script_match.start()].strip()

        points = []
        for line in content_wo_mermaid.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("#"):
                continue
            if stripped.startswith("- "):
                points.append(stripped[2:].strip())
            elif stripped.startswith("**") and stripped.endswith("**"):
                points.append(stripped.strip("*"))
            elif stripped.startswith("#### "):
                points.append(stripped[5:].strip())
            else:
                points.append(stripped)

        slides.append(SlideContent(title=title, points=points, script=script, mermaid=mermaid))
    return slides


def render_mermaid(mermaid_code: str, out_path: Path) -> None:
    encoded = base64.urlsafe_b64encode(mermaid_code.encode("utf-8")).decode("ascii")
    url = f"https://mermaid.ink/img/{encoded}?type=png&bgColor=ffffff"
    response = requests.get(url, timeout=60)
    response.raise_for_status()
    out_path.write_bytes(response.content)


def add_title(slide, title: str) -> None:
    shape = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12.0), Inches(0.65))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_subline(slide, text: str) -> None:
    shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.4))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY


def add_points_box(slide, points: List[str], top: float, height: float, title: str = "核心内容") -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(top),
        Inches(5.35),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = PURPLE
    box.line.width = Pt(1.2)

    title_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(top - 0.15),
        Inches(1.4),
        Inches(0.38),
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = PURPLE
    title_box.line.color.rgb = PURPLE
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(top + 0.25), Inches(4.9), Inches(height - 0.4))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)

    trimmed = points[:16]
    for idx, item in enumerate(trimmed):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(4)
        p.bullet = True
        r = p.runs[0]
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(14 if len(trimmed) <= 10 else 12)
        r.font.color.rgb = TEXT


def add_mermaid_image(slide, image_path: Path, top: float, height: float, title: str = "图示") -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.05),
        Inches(top),
        Inches(6.65),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = PURPLE
    box.line.width = Pt(1.2)

    title_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.3),
        Inches(top - 0.15),
        Inches(1.15),
        Inches(0.38),
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = BLUE
    title_box.line.color.rgb = BLUE
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    slide.shapes.add_picture(str(image_path), Inches(6.25), Inches(top + 0.25), width=Inches(6.25), height=Inches(height - 0.45))


def add_script_box(slide, script: str) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.72),
        Inches(12.0),
        Inches(0.55),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PURPLE_LIGHT
    box.line.color.rgb = PURPLE
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"讲解要点：{script[:220]}{'...' if len(script) > 220 else ''}"
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(10.5)
    r.font.color.rgb = TEXT


def add_gantt_chart(slide) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.05),
        Inches(1.45),
        Inches(6.65),
        Inches(4.9),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = PURPLE
    box.line.width = Pt(1.2)

    title_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.3),
        Inches(1.3),
        Inches(1.35),
        Inches(0.38),
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = BLUE
    title_box.line.color.rgb = BLUE
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "甘特图"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    months = ["5月", "6月", "7月", "8月", "9月", "10月"]
    tasks = [
        ("架构梳理与能力映射", 0, 1, PURPLE),
        ("权限与系统底座统一", 1, 2, RGBColor(164, 102, 204)),
        ("项目主数据中心建设", 1, 3, BLUE),
        ("自动化资源中心建设", 2, 4, RGBColor(55, 168, 116)),
        ("老项目功能拆分迁移", 3, 5, RGBColor(227, 132, 37)),
        ("流程打通与联调", 4, 5, RGBColor(227, 94, 94)),
        ("测试中心能力建设", 4, 5, RGBColor(112, 142, 67)),
        ("AI 场景设计与选型", 5, 6, RGBColor(66, 133, 244)),
        ("AI 助手一期接入", 5, 6, RGBColor(142, 68, 173)),
    ]

    start_x = 8.35
    month_width = 0.62
    top_y = 1.95

    for i, month in enumerate(months):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(start_x + month_width * i),
            Inches(top_y),
            Inches(month_width - 0.02),
            Inches(0.35),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PURPLE_LIGHT
        shape.line.color.rgb = PURPLE
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = month
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = TEXT

    for idx, (task, start, end, color) in enumerate(tasks):
        y = 2.38 + idx * 0.38
        label = slide.shapes.add_textbox(Inches(6.3), Inches(y), Inches(1.9), Inches(0.28))
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = task
        r = p.runs[0]
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(10)
        r.font.color.rgb = TEXT

        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(start_x + month_width * start),
            Inches(y),
            Inches(month_width * max(end - start, 0.7) - 0.04),
            Inches(0.24),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color


def build_presentation(slides: List[SlideContent]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[6]

    for idx, content in enumerate(slides, start=1):
        slide = prs.slides.add_slide(title_layout)
        add_title(slide, content.title)
        add_subline(slide, "SakurA Admin 老项目重构与 AI 自动化研发平台建设方案")

        if content.mermaid:
            add_points_box(slide, content.points, top=1.45, height=4.9)
            if "gantt" in content.mermaid:
                add_gantt_chart(slide)
            else:
                image_path = ASSETS / f"slide_{idx:02d}.png"
                try:
                    render_mermaid(content.mermaid, image_path)
                    add_mermaid_image(slide, image_path, top=1.45, height=4.9, title="架构图" if idx in {4, 5, 12} else "流程图")
                except Exception:
                    panel = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                        Inches(6.05),
                        Inches(1.45),
                        Inches(6.65),
                        Inches(4.9),
                    )
                    panel.fill.solid()
                    panel.fill.fore_color.rgb = PURPLE_LIGHT
                    panel.line.color.rgb = PURPLE
                    tf = panel.text_frame
                    p = tf.paragraphs[0]
                    p.text = "图示渲染失败，已保留文字内容"
                    p.alignment = PP_ALIGN.CENTER
                    r = p.runs[0]
                    r.font.name = "Microsoft YaHei"
                    r.font.size = Pt(20)
                    r.font.bold = True
                    r.font.color.rgb = PURPLE
        else:
            add_points_box(slide, content.points, top=1.45, height=5.4)
            # add a simple highlight panel on the right for visual balance
            panel = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(6.1),
                Inches(1.45),
                Inches(6.55),
                Inches(5.4),
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = PURPLE_LIGHT
            panel.line.color.rgb = PURPLE
            panel.line.width = Pt(1.2)
            tf = panel.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "本页重点"
            p.alignment = PP_ALIGN.CENTER
            r = p.runs[0]
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = PURPLE
            for item in content.points[:6]:
                p = tf.add_paragraph()
                p.text = item
                p.bullet = True
                p.space_before = Pt(8)
                r = p.runs[0]
                r.font.name = "Microsoft YaHei"
                r.font.size = Pt(16)
                r.font.color.rgb = TEXT

        if content.script:
            add_script_box(slide, content.script)
    return prs


def main() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    slides = parse_slides(SOURCE.read_text(encoding="utf-8"))
    prs = build_presentation(slides)
    prs.save(OUTPUT)
    print(f"Generated: {OUTPUT}")


if __name__ == "__main__":
    main()
